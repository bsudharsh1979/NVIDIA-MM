from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path

import nbformat

from .config import settings
from .db import (
    Course,
    Notebook,
    NotebookCell,
    SourceArtifact,
    SourceSpan,
    Session,
)
from .ids import artifact_uid, notebook_uid, span_uid

DANGEROUS = re.compile(
    r"\b(kubectl|helm|docker|rm\s+-rf|subprocess|os\.system|requests\.(get|post)|socket\.)\b",
    re.I,
)
HEADING = re.compile(r"^(#{1,6})\s+(.*)")
SHELL = re.compile(r"^!", re.M)


def _checksum(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()[:16]


def _hash_embed(text: str, dim: int = 64) -> list[float]:
    vec = [0.0] * dim
    tokens = re.findall(r"[a-z0-9\-\+_]+", text.lower())
    if not tokens:
        return vec
    for tok in tokens:
        h = int(hashlib.md5(tok.encode()).hexdigest(), 16)
        vec[h % dim] += 1.0
        vec[(h // dim) % dim] += 0.35
    norm = sum(v * v for v in vec) ** 0.5 or 1.0
    return [round(v / norm, 6) for v in vec]


def ingest_all(session: Session) -> dict:
    materials = settings.materials_path
    course = session.query(Course).filter_by(slug="nvidia-multimodal").one_or_none()
    if not course:
        course = Course(
            slug="nvidia-multimodal",
            title="NVIDIA — Building Multimodal AI Applications",
            description="Fusion, contrastive pre-training, projection, OCR, VSS, Graph-RAG, and CILP assessment.",
        )
        session.add(course)
        session.flush()

    counts = {"notebooks": 0, "pdf": 0, "pptx": 0, "cells": 0, "spans": 0}
    if not materials.exists():
        return counts

    for path in sorted(materials.rglob("*")):
        if path.suffix.lower() == ".ipynb":
            n, s = ingest_notebook(session, course.id, path)
            counts["notebooks"] += 1
            counts["cells"] += n
            counts["spans"] += s
        elif path.suffix.lower() == ".pdf":
            ingest_pdf(session, course.id, path)
            counts["pdf"] += 1
        elif path.suffix.lower() in {".ppt", ".pptx"}:
            ingest_pptx(session, course.id, path)
            counts["pptx"] += 1
        elif path.suffix.lower() in {".html", ".htm"}:
            ingest_html(session, course.id, path)
            counts["html"] = counts.get("html", 0) + 1
    session.commit()
    return counts


def ingest_notebook(session: Session, course_id: int, path: Path) -> tuple[int, int]:
    nb = nbformat.read(str(path), as_version=4)
    checksum = _checksum(path)
    existing = session.query(SourceArtifact).filter_by(filename=path.name, source_type="notebook").one_or_none()
    art_uid = artifact_uid(path.name)
    if existing and existing.checksum == checksum:
        if not getattr(existing, "uid", ""):
            existing.uid = art_uid
        nb_row = session.query(Notebook).filter_by(artifact_id=existing.id).one_or_none()
        if nb_row and not getattr(nb_row, "uid", ""):
            nb_row.uid = notebook_uid(path.name)
        return 0, 0
    if existing:
        session.delete(existing)
        session.flush()

    artifact = SourceArtifact(
        course_id=course_id,
        source_type="notebook",
        filename=path.name,
        title=_notebook_title(nb, path.name),
        checksum=checksum,
        extra={"cell_count": len(nb.cells)},
        uid=art_uid,
    )
    session.add(artifact)
    session.flush()

    order = _notebook_order(path.name)
    overview = NOTEBOOK_OVERVIEWS.get(path.name, {})
    notebook = Notebook(
        artifact_id=artifact.id,
        slug=path.stem,
        order_index=order,
        purpose=overview.get("purpose", artifact.title),
        why_it_matters=overview.get("why", ""),
        expected_outcome=overview.get("outcome", ""),
        uid=notebook_uid(path.name),
    )
    session.add(notebook)
    session.flush()

    cells_n = spans_n = 0
    current_heading = ""
    for idx, cell in enumerate(nb.cells):
        src = "".join(cell.source) if isinstance(cell.source, list) else str(cell.source or "")
        outputs = _serialize_outputs(cell)
        ctype = cell.cell_type
        heading = current_heading
        if ctype == "markdown":
            for line in src.splitlines():
                m = HEADING.match(line.strip())
                if m:
                    current_heading = m.group(2).strip()
                    heading = current_heading
                    break
        commands = []
        if ctype == "code":
            commands = [ln[1:] for ln in src.splitlines() if ln.strip().startswith("!")]
        dangerous = bool(DANGEROUS.search(src) or commands)
        extras = {
            "has_yaml": "apiVersion:" in src or src.strip().startswith("---"),
            "has_json": src.strip().startswith("{") or src.strip().startswith("["),
            "has_manifest": "kind:" in src and "apiVersion:" in src,
            "models": _find_models(src),
            "metrics": _find_metrics(src),
        }
        locator = {"source_type": "notebook", "file": path.name, "cell_index": idx}
        span = SourceSpan(
            artifact_id=artifact.id,
            span_type=ctype,
            locator=locator,
            title=heading or f"Cell {idx}",
            text=src if ctype == "markdown" else "",
            code=src if ctype == "code" else "",
            stored_output=outputs,
            heading=heading,
            embedding=_hash_embed(f"{path.name} {heading} {src}"[:8000]),
            uid=span_uid(art_uid, locator, ctype, idx),
        )
        session.add(span)
        session.flush()
        session.add(
            NotebookCell(
                notebook_id=notebook.id,
                span_id=span.id,
                cell_index=idx,
                cell_type=ctype,
                execution_count=cell.get("execution_count"),
                markdown=src if ctype == "markdown" else "",
                code=src if ctype == "code" else "",
                stored_output=outputs,
                heading=heading,
                commands=commands,
                dangerous=dangerous,
                extras=extras,
            )
        )
        cells_n += 1
        spans_n += 1
    return cells_n, spans_n


def ingest_pdf(session: Session, course_id: int, path: Path) -> None:
    try:
        import fitz
    except ImportError:
        return
    doc = fitz.open(path)
    artifact = SourceArtifact(
        course_id=course_id,
        source_type="pdf",
        filename=path.name,
        title=path.stem,
        checksum=_checksum(path),
        extra={"pages": doc.page_count},
    )
    session.add(artifact)
    session.flush()
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text")
        session.add(
            SourceSpan(
                artifact_id=artifact.id,
                span_type="page",
                locator={"source_type": "pdf", "file": path.name, "page": i},
                title=f"Page {i}",
                text=text,
                embedding=_hash_embed(text[:8000]),
            )
        )


def ingest_pptx(session: Session, course_id: int, path: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError:
        return
    prs = Presentation(str(path))
    artifact = SourceArtifact(
        course_id=course_id,
        source_type="pptx",
        filename=path.name,
        title=path.stem,
        checksum=_checksum(path),
        extra={"slides": len(prs.slides)},
    )
    session.add(artifact)
    session.flush()
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        notes = ""
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text_frame.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        body = "\n".join(texts)
        session.add(
            SourceSpan(
                artifact_id=artifact.id,
                span_type="slide",
                locator={"source_type": "pptx", "file": path.name, "slide": i},
                title=texts[0][:120] if texts else f"Slide {i}",
                text=body + (("\nNOTES:\n" + notes) if notes else ""),
                embedding=_hash_embed((body + notes)[:8000]),
            )
        )


def _serialize_outputs(cell) -> str:
    if cell.cell_type != "code":
        return ""
    chunks = []
    for out in cell.get("outputs") or []:
        otype = out.get("output_type")
        if otype == "stream":
            chunks.append("".join(out.get("text", [])))
        elif otype in {"execute_result", "display_data"}:
            data = out.get("data") or {}
            if "text/plain" in data:
                chunks.append("".join(data["text/plain"]) if isinstance(data["text/plain"], list) else str(data["text/plain"]))
        elif otype == "error":
            chunks.append("ERROR: " + "\n".join(out.get("traceback") or [])[:2000])
    return "\n".join(chunks)[:8000]


def _notebook_title(nb, fallback: str) -> str:
    for cell in nb.cells:
        if cell.cell_type == "markdown":
            src = "".join(cell.source) if isinstance(cell.source, list) else cell.source
            for line in src.splitlines():
                if line.startswith("# "):
                    return line[2:].strip()
    return fallback


def _notebook_order(name: str) -> int:
    order = [
        "00_jupyterlab.ipynb",
        "01a_Early_and_Late_Fusion.ipynb",
        "01b_Exploring_Modalities.ipynb",
        "02a_Intermediate_Fusion.ipynb",
        "02b_Contrastive_Pretraining.ipynb",
        "03a_Projection.ipynb",
        "03b_OCR_Pipelines.ipynb",
        "04a_VSS.ipynb",
        "04b_VSS_GraphRAG.ipynb",
        "05_Assessment.ipynb",
    ]
    return order.index(name) if name in order else 100 + int(artifact_uid(name)[:8], 16) % 50


def _find_models(src: str) -> list[str]:
    found = []
    for name in ("vila-1.5", "VGG16", "CLIP", "YOLOX", "nv-yolox", "CILP", "U-Net"):
        if name.lower() in src.lower():
            found.append(name)
    return found


def _find_metrics(src: str) -> list[str]:
    found = []
    for name in ("loss", "accuracy", "cosine", "chunk_duration", "temperature", "top_p"):
        if name in src:
            found.append(name)
    return found


def lexical_score(query: str, text: str) -> float:
    q = set(re.findall(r"[a-z0-9]+", query.lower()))
    t = re.findall(r"[a-z0-9]+", text.lower())
    if not q or not t:
        return 0.0
    tf = sum(1 for w in t if w in q)
    return tf / (len(t) ** 0.4)


def hybrid_search(session: Session, query: str, k: int = 8, filters: dict | None = None) -> list[dict]:
    qvec = _hash_embed(query)
    spans = session.query(SourceSpan).all()
    scored = []
    for span in spans:
        if filters and filters.get("file") and span.locator.get("file") != filters["file"]:
            continue
        emb = span.embedding or [0.0] * 64
        cosine = sum(a * b for a, b in zip(qvec, emb))
        lex = lexical_score(query, f"{span.title} {span.heading or ''} {span.text} {span.code} {span.locator}")
        if lex <= 0 and cosine < 0.2:
            continue
        score = 0.3 * cosine + 0.7 * min(lex, 2.0)
        blob = f"{span.title} {span.text} {span.code}".lower()
        for phrase in ("late fusion", "early fusion", "intermediate fusion", "colored cubes", "cilp", "chunk_duration", "graph-rag", "xyza"):
            if phrase in query.lower() and phrase.replace("-", " ") in blob.replace("-", " "):
                score += 0.5
        scored.append((score, span))
    scored.sort(key=lambda x: x[0], reverse=True)
    out = []
    for score, span in scored[:k]:
        out.append(
            {
                "score": round(float(score), 4),
                "title": span.title,
                "text": (span.text or span.code)[:1200],
                "locator": span.locator,
                "evidence_type": "COURSE_SOURCE",
            }
        )
    return out


NOTEBOOK_OVERVIEWS = {
    "00_jupyterlab.ipynb": {
        "purpose": "Orient to JupyterLab and GPU memory reset.",
        "why": "Labs share a GPU; leftover kernels poison later experiments.",
        "outcome": "Run a cell and reset the kernel without treating that as model training.",
    },
    "01a_Early_and_Late_Fusion.ipynb": {
        "purpose": "Compare RGB vs LiDAR vs early/late fusion on Omniverse synthetic objects.",
        "why": "Fusion choice depends on whether a modality carries the missing signal.",
        "outcome": "Explain XYZA conversion and why late fusion is an ensemble of heads.",
    },
    "01b_Exploring_Modalities.ipynb": {
        "purpose": "Turn audio into spectrograms and CT into slice animations.",
        "why": "New sensors become CNN/U-Net-shaped once represented as grids/volumes.",
        "outcome": "State Nyquist and why CT uses U-Net more than a 2D CNN.",
    },
    "02a_Intermediate_Fusion.ipynb": {
        "purpose": "On color-identical-shape cubes, compare early, late, concat, and matmul fusion.",
        "why": "When LiDAR cannot see color, unimodal LiDAR overfits; streams must mix mid-network.",
        "outcome": "Defend concat vs matmul vs late fusion for complementary sensors.",
    },
    "02b_Contrastive_Pretraining.ipynb": {
        "purpose": "Build a CLIP-style model pairing FashionMNIST photos with Sobel outlines.",
        "why": "Contrastive spaces let you retrieve or substitute a cheaper modality.",
        "outcome": "Write the all-pairs cosine matrix and symmetric CE loss.",
    },
    "03a_Projection.ipynb": {
        "purpose": "Project text embeddings into a frozen image classifier space.",
        "why": "You often cannot retrain a production unimodal model; you add a projector.",
        "outcome": "Wire projector output into get_img_embs-sized features.",
    },
    "03b_OCR_Pipelines.ipynb": {
        "purpose": "Extract text, tables, and figures from the GB200 NVL72 datasheet.",
        "why": "LLM RAG over PDFs fails if chunking and layout detection are naive.",
        "outcome": "Contrast naive vs by_title chunking and NV-YOLOX labels.",
    },
    "04a_VSS.ipynb": {
        "purpose": "Drive NVIDIA VSS summarization APIs and CA-RAG.",
        "why": "Video understanding is a VLM+LLM pipeline, not a single prompt.",
        "outcome": "Relate chunk_duration to processed frames and prompt persona to caption quality.",
    },
    "04b_VSS_GraphRAG.ipynb": {
        "purpose": "Q&A with Vector-RAG vs Graph-RAG; inspect Neo4j.",
        "why": "Relational questions (who wore PPE) need edges, not only similar captions.",
        "outcome": "Name G-Extraction, G-Retriever, G-Generation and enable_chat.",
    },
    "05_Assessment.ipynb": {
        "purpose": "Pass CILP + projector so RGB images reuse a frozen LiDAR classifier.",
        "why": "Cameras are cheaper than LiDAR; contrastive alignment transfers the head.",
        "outcome": "Hit loss and accuracy gates without unfreezing lidar_cnn.",
    },
}


def ingest_html(session: Session, course_id: int, path: Path) -> None:
    raw = path.read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<script[\s\S]*?</script>", " ", raw, flags=re.I)
    text = re.sub(r"<style[\s\S]*?</style>", " ", text, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    uid = artifact_uid(path.name)
    artifact = SourceArtifact(
        course_id=course_id,
        source_type="html",
        filename=path.name,
        title=path.stem,
        checksum=_checksum(path),
        extra={"chars": len(text)},
        uid=uid,
    )
    session.add(artifact)
    session.flush()
    locator = {"source_type": "html", "file": path.name, "page": 1}
    session.add(
        SourceSpan(
            artifact_id=artifact.id,
            span_type="page",
            locator=locator,
            title=path.stem,
            text=text[:20000],
            embedding=_hash_embed(text[:8000]),
            uid=span_uid(uid, locator, "page", 1),
        )
    )


CELL_TEACHING = {
    "why": "This cell exists to make a course idea executable or to record a result the learner should predict first.",
    "verify": "Read stored outputs if present. Do not assume a blank output means the command succeeded.",
    "failure": "GPU OOM, missing data/ paths, and network calls to via-server/ngc-client will fail outside the DLI classroom.",
    "modify": "Change a hyperparameter or prompt on paper first, then predict the qualitative effect before any real run.",
    "business": "Operational cost, safety Q&A, or exam pass/fail depends on getting this cell's idea right — not on decorating the notebook.",
}
